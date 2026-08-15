"""Insere no deck os slides de ANÁLISE dos experimentos do MLflow.

Acrescenta dois slides logo depois do slide 20 (o de rastreamento MLflow):

  A. "Conversão não é lucro"  — o gráfico valor × conversão + a leitura
  B. "Comparação e robustez"  — painel de métricas + sensibilidade à seed

Pré-requisito: rode antes `python scripts\\build_mlflow_charts.py`, que gera os
PNGs em `docs/img/mlflow/` a partir dos runs reais.

O script é **idempotente**: se os slides já existirem (detectados pelo título),
ele os remove antes de recriar — então dá para rodar de novo depois de regerar
os gráficos, sem duplicar nada.

    python scripts\\add_mlflow_slides.py
"""

from __future__ import annotations

import shutil
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
DECK = PROJECT_ROOT / "docs" / "Adaptive-Offers-Pitch-Grupo74.pptx"
IMG = PROJECT_ROOT / "docs" / "img" / "mlflow"

# Paleta do deck
PANEL = RGBColor(0x03, 0x0D, 0x24)
PANEL2 = RGBColor(0x01, 0x0A, 0x1A)
CYAN = RGBColor(0x1A, 0x6F, 0xFF)
GREEN = RGBColor(0x1A, 0x9E, 0x1A)
GOLD = RGBColor(0xFF, 0xC2, 0x00)
RED = RGBColor(0xE8, 0x40, 0x00)
TEXT = RGBColor(0xED, 0xED, 0xED)
MUTED = RGBColor(0x88, 0x99, 0xBB)

TITLE_A = "MLflow — Conversão não é lucro"
TITLE_B = "MLflow — Comparação e robustez"
ANCHOR = "MLflow — Rastreamento de Experimentos"

# Títulos que estes slides já tiveram. A idempotência compara por título, então
# uma versão anterior sem acento não seria reconhecida — e o script duplicaria
# os slides em vez de substituí-los.
LEGACY_TITLES = {
    "MLflow — Conversao nao e lucro",
    "MLflow — Comparacao e robustez",
}
OWN_TITLES = {TITLE_A, TITLE_B} | LEGACY_TITLES


# ── helpers de desenho (mesmo estilo de scripts/fix_pptx_slides.py) ──────────
def clear_placeholders(slide) -> None:
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)


def rect(slide, left, top, w, h, fill, border=None, bw: float = 1.0):
    s = slide.shapes.add_shape(1, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s


def txt(slide, left, top, w, h, text, size=11, color=TEXT, bold=False,
        align=PP_ALIGN.LEFT, wrap=True, name="Calibri", space_after=6,
        line_spacing=None):
    """Cria uma caixa de texto com a formatacao aplicada NO RUN.

    Cuidado que custou caro: em python-pptx, `paragraph.font` escreve em
    `a:pPr/a:defRPr` — que e so um DEFAULT para runs sem formatacao propria. O
    PowerPoint respeita esse default, mas visualizadores de PDF, LibreOffice e
    Google Slides nao: eles caem na cor do tema (escura, neste deck) e no
    tamanho padrao, deixando o texto ilegivel e quebrando o wrap. A formatacao
    precisa ir em `run.font`, que gera `a:rPr` no proprio run.
    """
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0

    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        f = run.font                    # <- run-level: vira a:rPr
        f.size = Pt(size)
        f.color.rgb = color
        f.bold = bold
        f.name = name
    return box


def header(slide, prs, title, subtitle=""):
    w = prs.slide_width
    rect(slide, 0, 0, w, Inches(0.06), GOLD)
    rect(slide, 0, Inches(0.06), w, Inches(0.74), PANEL2)
    txt(slide, Inches(0.45), Inches(0.08), w - Inches(0.9), Inches(0.46),
        title, size=24, bold=True, color=TEXT, wrap=False)
    if subtitle:
        txt(slide, Inches(0.45), Inches(0.54), w - Inches(0.9), Inches(0.24),
            subtitle, size=11, color=MUTED, wrap=False)


def footer(slide, prs, text):
    w, h = prs.slide_width, prs.slide_height
    txt(slide, Inches(0.45), h - Inches(0.45), w - Inches(0.9), Inches(0.3),
        text, size=9, color=MUTED, wrap=False)


def move_slide(prs, from_idx: int, to_idx: int) -> None:
    lst = prs.slides._sldIdLst
    el = lst[from_idx]
    lst.remove(el)
    lst.insert(to_idx, el)


def drop_slide(prs, idx: int) -> None:
    lst = prs.slides._sldIdLst
    el = lst[idx]
    prs.part.drop_rel(el.get(qn("r:id")))
    lst.remove(el)


def slide_title(slide) -> str:
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip()
    return ""


# ── os dois slides ──────────────────────────────────────────────────────────
def build_slide_a(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    w, h = prs.slide_width, prs.slide_height
    clear_placeholders(slide)
    rect(slide, 0, 0, w, h, PANEL)

    header(slide, prs, TITLE_A,
           "O experimento que sustenta a tese do pitch — leitura direta dos runs")

    img = IMG / "mlflow-valor-vs-conversao.png"
    slide.shapes.add_picture(str(img), Inches(0.40), Inches(1.05),
                             width=Inches(7.85))

    # Painel de leitura à direita.
    # Sem quebras de linha manuais: o wrap resolve. Linha fixa "\n" so acerta se
    # a fonte renderizar exatamente com a metrica prevista — em outro
    # visualizador ela quebra no meio das palavras.
    px, pw = Inches(8.50), Inches(4.40)
    tw = pw - Inches(0.50)                      # largura util do texto
    tx = px + Inches(0.25)
    rect(slide, px, Inches(1.05), pw, Inches(5.35), PANEL2, border=GOLD, bw=1.5)

    txt(slide, tx, Inches(1.24), tw, Inches(0.26),
        "A LEITURA", size=10, bold=True, color=GOLD)

    txt(slide, tx, Inches(1.60), tw, Inches(1.05),
        "O Baseline tem a MAIOR conversão (10,6%) e o MENOR valor (R$ 76.560).",
        size=15, bold=True, color=TEXT, line_spacing=1.15)

    txt(slide, tx, Inches(2.78), tw, Inches(1.20),
        "Ele empurra quase tudo para a oferta de maior margem e converte em "
        "volume — mas ignora o contexto e queima as decisões de maior valor "
        "esperado.",
        size=12, color=TEXT, line_spacing=1.2)

    txt(slide, tx, Inches(4.05), tw, Inches(1.05),
        "O LinUCB converte MENOS (8,6%) e entrega +41% de valor: escolhe a "
        "oferta certa para o perfil certo, não a mais fácil de vender.",
        size=12, color=TEXT, line_spacing=1.2)

    rect(slide, tx, Inches(5.22), tw, Inches(1.00), PANEL, border=CYAN, bw=1.2)
    txt(slide, tx + Inches(0.16), Inches(5.38), tw - Inches(0.32), Inches(0.70),
        "Por isso a recompensa é ponderada por margem: "
        "value = P(conversão) × margem.",
        size=12, bold=True, color=CYAN, line_spacing=1.2)

    footer(slide, prs,
           "Fonte: runs do MLflow (train-all, 6.000 rounds, seed=42)  |  "
           "reproduza com: python scripts/mlflow_report.py --compare")
    return slide


def build_slide_b(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    w, h = prs.slide_width, prs.slide_height
    clear_placeholders(slide)
    rect(slide, 0, 0, w, h, PANEL)

    header(slide, prs, TITLE_B,
           "As quatro métricas que a banca olha — e por que 1 seed não decide nada")

    slide.shapes.add_picture(str(IMG / "mlflow-metricas.png"),
                             Inches(0.35), Inches(1.02), width=Inches(6.30))
    slide.shapes.add_picture(str(IMG / "mlflow-sensibilidade-seed.png"),
                             Inches(6.80), Inches(1.02), width=Inches(6.20))

    # Faixa de leitura no rodape (sem quebras manuais — ver nota em txt())
    by = Inches(5.50)
    rect(slide, Inches(0.35), by, w - Inches(0.70), Inches(1.32), PANEL2,
         border=GOLD, bw=1.5)
    txt(slide, Inches(0.58), by + Inches(0.13), w - Inches(1.16), Inches(0.26),
        "O QUE ISSO SIGNIFICA", size=10, bold=True, color=GOLD)
    txt(slide, Inches(0.58), by + Inches(0.47), Inches(5.75), Inches(0.72),
        "LinUCB tem o MENOR regret (9,7%) — e o baseline, sem exploração, "
        "o maior (36,7%). Explorar ~20% do tráfego se paga.",
        size=12, color=TEXT, line_spacing=1.2)
    txt(slide, Inches(6.85), by + Inches(0.47), Inches(5.75), Inches(0.72),
        "Trocando só a seed, o baseline varia +37%. Por isso a conclusão do "
        "relatório usa 5 seeds: LinUCB lidera na média, com CV de 2,97%.",
        size=12, color=TEXT, line_spacing=1.2)

    footer(slide, prs,
           "Esquerda: MLflow, seed=42  |  Direita: seed=42 (train-all) vs seed=123 "
           "(evaluate — a fonte do slide 09)  |  6.000 rounds")
    return slide


def main() -> None:
    if not DECK.exists():
        raise SystemExit(f"deck nao encontrado: {DECK}")
    faltando = [p.name for p in (
        IMG / "mlflow-valor-vs-conversao.png",
        IMG / "mlflow-metricas.png",
        IMG / "mlflow-sensibilidade-seed.png",
    ) if not p.exists()]
    if faltando:
        raise SystemExit(
            f"PNG(s) ausente(s): {', '.join(faltando)}\n"
            "rode antes: python scripts\\build_mlflow_charts.py")

    backup = DECK.with_suffix(f".pptx.bak-{datetime.now():%Y%m%d-%H%M%S}")
    shutil.copy2(DECK, backup)
    print(f"backup: {backup.name}")

    # ── Passe 1: remover versões anteriores destes slides ────────────────────
    # Em DOIS passes de propósito. O python-pptx tira o sldId da lista mas deixa
    # a *part* no pacote; se removermos e adicionarmos slides na mesma sessão, os
    # nomes de part colidem ("Duplicate name: ppt/slides/slideNN.xml") e um slide
    # existente é sobrescrito silenciosamente. Salvar e reabrir entre os passes
    # faz o pacote ser reserializado sem as parts órfãs.
    prs = Presentation(DECK)
    removidos = 0
    for i in range(len(prs.slides) - 1, -1, -1):
        antigo = slide_title(prs.slides[i])
        if antigo in OWN_TITLES:
            drop_slide(prs, i)
            removidos += 1
            print(f"  slide antigo removido (S{i + 1}): {antigo}")
    if removidos:
        prs.save(DECK)
        prs = Presentation(DECK)

    # ── Passe 2: inserir ─────────────────────────────────────────────────────
    # Ancora: logo apos o slide de rastreamento MLflow
    anchor = next((i for i, s in enumerate(prs.slides)
                   if slide_title(s) == ANCHOR), None)
    if anchor is None:
        anchor = len(prs.slides) - 2
        print(f"  [aviso] slide '{ANCHOR}' nao encontrado; inserindo em {anchor + 1}")

    build_slide_a(prs)
    move_slide(prs, len(prs.slides) - 1, anchor + 1)
    build_slide_b(prs)
    move_slide(prs, len(prs.slides) - 1, anchor + 2)

    prs.save(DECK)
    print(f"\n2 slides inseridos apos o slide {anchor + 1}. Total: {len(prs.slides)} slides.")
    for i, s in enumerate(prs.slides, 1):
        if slide_title(s) in (TITLE_A, TITLE_B, ANCHOR):
            print(f"   S{i:02d}  {slide_title(s)}")


if __name__ == "__main__":
    main()
