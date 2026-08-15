"""Sincroniza o deck com os números reais do projeto (base UCI, não o fac-símile).

O deck foi montado quando o pipeline ainda rodava sobre o fac-símile determinístico
e carregava números inflados (+66,6% de lift, R$ 424.820 de reward). Depois que a
base real do Kaggle entrou (`provenance="real"`), `reports/technical-report.md` e
`artifacts/evaluation_report.json` passaram a mostrar ganhos *single digit* — mas o
deck ficou para trás, e os slides 09 e 20 passaram a se contradizer.

Este script reescreve **só o texto** dos runs afetados (preservando fonte, cor e
posição) para que o deck conte a mesma história dos relatórios.

Uso:  python scripts\\sync_pptx_numbers.py
"""

from __future__ import annotations

import shutil
from datetime import datetime
from pathlib import Path

from pptx import Presentation

DECK = Path("docs/Adaptive-Offers-Pitch-Grupo74.pptx")

# Shapes com várias linhas: (slide, shape, parágrafo, texto_antigo, texto_novo)
PARAGRAPH_EDITS: list[tuple[int, int, int, str, str]] = [
    # --- S14 · FinOps citava o lift do fac-símile -----------------------------
    (
        13,
        7,
        0,
        "ROI: o ganho de valor (+60%) dilui o custo de compute/LLM",
        "ROI: o ganho de valor (+8–9% na base real) dilui o custo de compute/LLM",
    ),
]

# (slide_index_0based, shape_index, texto_antigo_esperado, texto_novo)
EDITS: list[tuple[int, int, str, str]] = [
    # --- S10 · portas erradas na tela de demonstração -------------------------
    (9, 15, "Streamlit · :8510", "Streamlit · :8503"),
    (9, 20, "mlflow ui · :5000", "mlflow ui · :5001"),

    # --- S13 · contagem de testes --------------------------------------------
    (12, 26, "61 testes (unit + integração)", "93 testes (unit + integração)"),

    # --- S20 · o que o MLflow realmente registra ------------------------------
    (19, 12, "  artefatos: modelo + features", "  tags: policy, stage, seed"),
    (19, 16, "  LinUCB promovido a Production", "  LinUCB e a politica ativa (v1)"),
    (19, 17, "  versoes anteriores em Staging", "  registry em artifacts/policies"),
    (19, 18, "  rollback disponivel", "  rollback em 1 comando"),
    (19, 19, "  API usa versao Production", "  API sempre le a versao ativa"),

    # --- S20 · resultados-chave (eram do fac-símile) --------------------------
    (19, 23, "  LinUCB: regret ratio 5,1%", "  LinUCB: regret ratio 8,3%"),
    (19, 24, "  Lift +66,6% vs baseline", "  Lift +8,2% vs baseline"),
    (19, 25, "  Convergencia: round ~800", "  DR-OPE 19,18 [17,66-20,65]"),

    # --- S20 · tabela, na mesma ordem do slide 09 (reward decrescente) --------
    # linha 1
    (19, 38, "LinUCB (campeao)", "Thompson"),
    (19, 40, "R$ 424.820", "R$ 114.290"),
    (19, 42, "5,1%", "11,8%"),
    (19, 44, "9,1%", "7,1%"),
    (19, 46, "+66,6%", "+9,2%"),
    # linha 2
    (19, 48, "Thompson", "LinUCB (recomendada)"),
    (19, 50, "R$ 351.200", "R$ 113.230"),
    (19, 52, "17,4%", "8,3%"),
    (19, 54, "7,8%", "9,1%"),
    (19, 56, "+37,9%", "+8,2%"),
    # linha 3
    (19, 58, "Nilos-UCB", "Baseline (greedy)"),
    (19, 60, "R$ 330.100", "R$ 104.700"),
    (19, 62, "22,5%", "10,9%"),
    (19, 64, "7,3%", "6,2%"),
    (19, 66, "+29,9%", "---"),
    # linha 4
    (19, 68, "Baseline (greedy)", "Nilos-UCB"),
    (19, 70, "R$ 254.990", "R$ 102.020"),
    (19, 72, "100%", "17,0%"),
    (19, 74, "5,5%", "7,1%"),
    (19, 76, "---", "-2,6%"),
    # --- S21 · a API cresceu de 7 para 14 endpoints desde o deck --------------
    (
        20,
        4,
        "FastAPI · 7 endpoints · auditavel · localhost:8000/docs",
        "FastAPI · 14 endpoints · auditavel · localhost:8000/docs",
    ),
    (
        20,
        43,
        "Swagger UI em http://localhost:8000/docs  |  Decisoes auditaveis em audit_log.jsonl"
        "  |  FIAP 7MLET Grupo 74",
        "Swagger UI em http://localhost:8000/docs  |  Decisoes auditaveis em"
        " artifacts/decisions/audit.jsonl  |  FIAP 7MLET Grupo 74",
    ),

    # rodapé
    (
        19,
        77,
        "LinUCB eleito campeao via MLflow Model Registry  |  Simulacao em 6.000 rounds"
        "  |  Bank Marketing Dataset  |  41.188 clientes reais",
        "LinUCB recomendado: menor regret, mais estavel em 5 seeds e aprovado no gate"
        " DR-OPE  |  6.000 rounds  |  Bank Marketing (UCI)  |  41.188 contatos reais",
    ),
]


def set_text(shape, new_text: str) -> None:
    """Troca o texto mantendo a formatação do primeiro run e removendo os demais."""
    paragraph = shape.text_frame.paragraphs[0]
    runs = paragraph.runs
    if not runs:
        shape.text_frame.text = new_text
        return
    runs[0].text = new_text
    for extra in runs[1:]:
        extra._r.getparent().remove(extra._r)


def main() -> None:
    if not DECK.exists():
        raise SystemExit(f"deck não encontrado: {DECK}")

    backup = DECK.with_suffix(f".pptx.bak-{datetime.now():%Y%m%d-%H%M%S}")
    shutil.copy2(DECK, backup)
    print(f"backup: {backup.name}")

    prs = Presentation(DECK)
    applied, skipped = 0, []
    total = len(EDITS) + len(PARAGRAPH_EDITS)

    for slide_idx, shape_idx, expected, new in EDITS:
        shape = prs.slides[slide_idx].shapes[shape_idx]
        current = shape.text_frame.text
        if current != expected:
            skipped.append(f"S{slide_idx + 1:02d}[{shape_idx}] esperava {expected!r}, achou {current!r}")
            continue
        set_text(shape, new)
        applied += 1

    for slide_idx, shape_idx, para_idx, expected, new in PARAGRAPH_EDITS:
        shape = prs.slides[slide_idx].shapes[shape_idx]
        paragraph = shape.text_frame.paragraphs[para_idx]
        if paragraph.text != expected:
            skipped.append(
                f"S{slide_idx + 1:02d}[{shape_idx}]p{para_idx} esperava {expected!r}, "
                f"achou {paragraph.text!r}"
            )
            continue
        runs = paragraph.runs
        runs[0].text = new
        for extra in runs[1:]:
            extra._r.getparent().remove(extra._r)
        applied += 1

    if skipped:
        print(f"\n{len(skipped)} edição(ões) NÃO aplicada(s) — o deck mudou desde o script:")
        for item in skipped:
            print(f"  - {item}")

    prs.save(DECK)
    print(f"\n{applied}/{total} edições aplicadas em {DECK}")


if __name__ == "__main__":
    main()
